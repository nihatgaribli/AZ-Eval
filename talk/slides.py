"""CMAI 2026 çıxışının slaydları - bir mənbədən iki format.

    python talk/slides.py            # hər ikisi
    python talk/slides.py --pptx     # yalnız PowerPoint
    python talk/slides.py --beamer   # yalnız LaTeX

NİYƏ BİR SKRİPT: iki dest ayrıca redaktə olunsa, məşq zamanı biri dəyişir,
digəri qalır və çıxış günü hansının cari olduğu bilinmir. Məzmun burada BİR
dəfə yazılır, hər iki fayl ondan qurulur.

DİL: slaydlar İNGİLİSCƏDİR - konfransın işçi dili ingiliscədir. Azərbaycan dili
yalnız nümunə suallarda görünür.

BEAMER ÜÇÜN XELATEX LAZIMDIR: nümunələrdə `ə` (U+0259) hərfi var və standart
`pdflatex` T1 şriftlərində bu simvol YOXDUR - sətir səssizcə pozulur. `xelatex`
Unicode şrift işlədir və problem aradan qalxır.
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path

HERE = Path(__file__).resolve().parent

#: Şrift. Windows 11-də hər zaman mövcuddur və `ə` daxil olmaqla Azərbaycan
#: latın əlifbasını tam örtür. Hər iki formatda eyni şrift işlənir.
FONT = "Segoe UI"
BEAMER_FONT = FONT


#: Vizual kimlik. Rənglər `talk/figures.py`-dakı təsdiqlənmiş palitradan gəlir,
#: yəni slayd və qrafik eyni mavini işlədir. Ayrı-ayrı seçilsəydi, proyektorda
#: iki fərqli mavi görünərdi.
ACCENT = "2a78d6"
INK = "0b0b0b"
INK_SOFT = "52514e"
MUTED = "898781"

#: Fon. QƏSDƏN xalis ağ (ffffff) deyil.
#:
#: Xalis ağ proyektorda kəskin görünür və qrafiklərin fonu ilə uyuşmur, çünki
#: qrafiklər `#fcfcfb` səth üzərində yoxlanılıb. Eyni ton həm slaydda, həm
#: qrafikdə işlənəndə şəkil slaydın içində "üzmür", ona yapışır.
SURFACE = "fcfcfb"

#: Sol kənar boşluğu. Hər element eyni xəttdən başlayır: başlıq, vurğu zolağı,
#: cədvəl, mətn, altlıq. Vahid şaquli xətt slaydlara səliqə verir və bunu
#: dinləyici şüurlu görmür, sadəcə hiss edir.
_MARGIN_INCHES = 0.7


@dataclass
class Slide:
    title: str
    #: Madde-madde mətn. Boş sətir aralıq buraxır.
    bullets: list[str] = field(default_factory=list)
    #: (başlıq sətri, sətirlər) sadə cədvəl.
    table: tuple[list[str], list[list[str]]] | None = None
    #: Slaydın altındakı bir cümləlik vurğu.
    punch: str = ""
    #: Qrafik faylının adı (`talk/` qovluğunda). Cədvəli əvəz edir.
    image: str = ""
    #: Danışıq qeydi. PPTX-də notlar bölməsinə, Beamer-də `\note{}`-a düşür.
    notes: str = ""


TITLE = "AZ-Eval: A Parallel Azerbaijani–English Benchmark Reveals\nNegative Cross-Lingual Transfer from Kazakh Fine-Tuning"
AUTHOR = "Nihat Garibli"
AFFILIATION = "Baku State University, Azerbaijan"
VENUE = "CMAI 2026 · Nazarbayev University, Astana"

MAIN: list[Slide] = [
    Slide(
        title="Open multilingual models are measured on high-resource languages",
        bullets=[
            "Kazakh has open evaluation sets - freshqa_kazakh, defan_kazakh, BeyneleBench",
            "Azerbaijani, to our knowledge, has none",
            "",
            "Both are Turkic. Kazakh is written in Cyrillic, Azerbaijani in Latin.",
        ],
        notes=(
            "One sentence only: this is the host institute's own methodology, "
            "extended to another Turkic language. Do not linger."
        ),
    ),
    Slide(
        title="The question",
        punch="Does fine-tuning on Kazakh transfer to Azerbaijani?",
        bullets=[
            "",
            "Kazakh is the closest well-resourced Turkic relative of Azerbaijani.",
            "Adapting to a related language is the intuitive move for a low-resource language.",
        ],
        notes="Stress the script difference now. It detonates on slide 8.",
    ),
    Slide(
        title="AZ-Eval - an open parallel benchmark",
        table=(
            ["", ""],
            [
                ["356 parallel AZ/EN short-answer items", "255 hand-written by a native speaker"],
                ["Every item human-verified", "101 from templated Wikidata queries"],
                ["7 categories", "math, history, geography, science, culture, language, world"],
                ["Accepted answers expanded", "rule-based Azerbaijani morphology"],
            ],
        ),
        punch="verified_by is a gate, not metadata: unverified rows never enter the set.",
        notes=(
            "Say in one sentence why the questions are not about Azerbaijan: we need "
            "facts the model demonstrably knows in English, otherwise the gap collapses "
            "to zero and measures nothing. Otherwise it comes back in Q&A."
        ),
    ),
    Slide(
        title="Four nested normalizations",
        punch="STRICT  →  MORPH  →  LENIENT  →  TRANSLIT",
        bullets=[
            "",
            "suffix stripping   ·   diacritic folding   ·   Cyrillic→Latin",
            "",
            "Each stage adds exactly one transformation, applied symmetrically to",
            "prediction and reference - so each score increment attributes error mass",
            "to that single surface phenomenon.",
            "",
            "Paired comparisons · 1000-resample bootstrap · permutation test · Holm correction",
        ],
        notes=(
            "This is a mathematics conference. Showing bare percentages looks weak; "
            "showing that the measurement is trustworthy looks strong. Do not claim a theorem."
        ),
    ),
    Slide(
        title="Result 1 - the gap is universal",
        table=(
            ["Model", "AZ", "EN", "Gap", "p (Holm)"],
            [
                ["Qwen3-1.7B", "5.9%", "23.6%", "17.7", "0.0020"],
                ["Qwen3-VL-4B-Instruct", "16.6%", "30.6%", "14.0", "0.0020"],
                ["Qwen3-VL-4B-Thinking", "12.6%", "29.5%", "16.9", "0.0020"],
                ["issai/Qolda-AVL-5B", "3.1%", "29.5%", "26.4", "0.0020"],
            ],
        ),
        punch="Majority-class baseline: 1.1%",
        notes=(
            "'Every model is bad at Azerbaijani' is the expected result - do not dwell. "
            "Point at the last row and move on: something is wrong there."
        ),
    ),
    Slide(
        title="Result 2: transfer is NEGATIVE",
        image="fig_negative_transfer.png",
        punch="Holm-corrected p = 0.0040   ·   its own declared base, same quantization, same prompt",
        notes=(
            "Stop here. This is the centre of the talk. Do not rush. "
            "The only difference between the two models is the Kazakh fine-tune."
        ),
    ),
    Slide(
        title="The mechanism is orthographic",
        image="fig_script_counts.png",
        punch="Understands the question, finds the answer, writes it in the wrong alphabet.",
        notes=(
            "The memorable moment of the talk. Slow down. Two sentences: 326 vs 0 is a "
            "direct count, it depends on no statistical test. And demanding the Latin "
            "alphabet in the prompt moves only 12 of the 326."
        ),
    ),
    Slide(
        title="Is it simply a weaker model?",
        bullets=[
            "96 world + science items - facts every model demonstrably knows in English",
        ],
        table=(
            ["Language", "Base model", "Kazakh-tuned", "Gap", "p (Holm)"],
            [
                ["English", "62.5%", "59.4%", "3.1", "1.0000"],
                ["Azerbaijani (STRICT)", "37.5%", "10.4%", "27.1", "0.0050"],
                ["Azerbaijani (TRANSLIT)", "38.5%", "33.3%", "5.2", "1.0000"],
            ],
        ),
        punch="Indistinguishable in English. 27.1 points apart in Azerbaijani - and transliteration closes it.",
        notes=(
            "A weaker model would be weaker in both languages. This one is weaker "
            "only in Azerbaijani. Raise the objection before the audience does."
        ),
    ),
    Slide(
        title="Takeaway",
        bullets=[
            "1.  Fine-tuning on a closely related language can HARM rather than help",
            "     - when the writing systems differ.",
            "",
            "2.  Most of the damage is orthographic, not semantic",
            "     - transliteration removes 56% of the deficit.",
            "",
            "3.  Which means it is in principle fixable",
            "     - e.g. by constraining the output script at decoding time.",
        ],
        punch="Dataset CC BY 4.0 · Code MIT · github.com/nihatgaribli/AZ-Eval",
        notes="End on the practical sentence. Someone building a low-resource model is in this room.",
    ),
]

BACKUP: list[Slide] = [
    Slide(
        title="Backup - where the gap is widest",
        table=(
            ["Category", "n", "Qwen3-VL-4B (AZ / EN)", "Qolda-AVL-5B (AZ / EN)"],
            [
                ["mathematics", "28", "3.6% / 25.0%", "3.6% / 28.6%"],
                ["world", "45", "51.1% / 57.8%", "0.0% / 60.0%"],
                ["history", "45", "4.4% / 0.0%", "0.0% / 2.2%"],
            ],
        ),
        punch="Mathematics: knowledge is present in English, the language blocks it.",
        notes=(
            "Keep this ready - it is a mathematics conference and this is the widest gap. "
            "History sits near zero in BOTH languages: that is ignorance, not a script problem."
        ),
    ),
    Slide(
        title="Backup - what survives transliteration",
        bullets=[
            "Transliteration removes 56% of the deficit - but not all of it.",
            "",
            "Residual after TRANSLIT: 5.9 points, Holm p = 0.0364 - still significant.",
            "",
            "So the claim is: the damage is LARGELY orthographic, not purely orthographic.",
        ],
        notes="Honest framing. At n=96 this residual was not detectable; at n=356 it is.",
    ),
    Slide(
        title="Backup - how the dataset was built",
        bullets=[
            "255 items hand-written by a native Azerbaijani speaker",
            "101 from 8 templated SPARQL queries × 4 syntactic variants (round-robin)",
            "",
            "No LLM authored, translated, or answered any dataset item.",
            "Rejected candidates are retained, so the acceptance rate stays auditable.",
            "",
            "Items whose gold answer depends on a contested position are excluded (6 removed).",
        ],
    ),
    Slide(
        title="Backup - limitations",
        bullets=[
            "n = 356. Confidence intervals are ±2–5 points; every reported comparison",
            "survives Holm correction.",
            "",
            "One model pair. The 9B/8B pair did not fit in 8 GB of VRAM.",
            "",
            "Two knowledge regimes are mixed on purpose - and kept in separate categories",
            "so the distinction stays visible instead of averaging away.",
        ],
    ),
    Slide(
        title="Backup - answer extraction",
        bullets=[
            "Identical rules for both languages:",
            "",
            "first line only  ·  strip 'Cavab:' / 'Answer:' prefixes",
            "strip markdown emphasis and trailing punctuation  ·  remove thinking blocks",
            "",
            "Raw model outputs are kept untouched (8 runs), so the scoring rule can be",
            "changed and the analysis re-run without spending another GPU-hour.",
        ],
    ),
]


# --------------------------------------------------------------------------
# PowerPoint
# --------------------------------------------------------------------------


def _rgb(hex_value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_value)


def _margin():
    """Sol kənar. `pptx` opsional asılılıqdır, ona görə modul səviyyəsində
    deyil, çağırış anında qurulur."""
    from pptx.util import Inches

    return Inches(_MARGIN_INCHES)


def _text(box, lines, size, color=INK, bold=False, space_after=8):
    """Mətn qutusunu doldurur və hər sətrin şriftini təyin edir.

    Boş sətir aralıq buraxır. python-pptx boş paraqrafa `run` yaratmır, ona
    görə şrift təyini yalnız run varsa aparılır.
    """
    from pptx.util import Pt

    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(space_after)
        if not paragraph.runs:
            continue
        run = paragraph.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = FONT
        run.font.color.rgb = _rgb(color)


def _paint(slide, prs) -> None:
    """Slaydın fonunu isti ağa boyayır.

    Xalis ağ (`ffffff`) proyektorda kəskin görünür və qrafiklərin fonu ilə
    uyuşmur: qrafiklər `#fcfcfb` səth üzərində yoxlanılıb. Eyni ton işlənəndə
    şəkil slaydın içində üzmür, ona yapışır.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(SURFACE)
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Fon birinci çəkilir, amma python-pptx onu sona əlavə edir, ona görə
    # z-sırasında ƏN ARXAYA köçürülür. Köçürülməsə bütün məzmunu örtür.
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _rule(slide, left, top, width, height_pt=3.5, color=ACCENT) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(height_pt))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(color)
    bar.line.fill.background()
    bar.shadow.inherit = False


def _footer(slide, number: int) -> None:
    """Slayd nömrəsi və qısa başlıq.

    Zalda kimsə "8-ci slaydda nə yazmışdın?" deyə soruşanda nömrə lazım olur.
    Susqun rəngdədir, çünki məlumat deyil, naviqasiyadır.
    """
    from pptx.util import Inches

    box = slide.shapes.add_textbox(_margin(), Inches(6.95), Inches(11.9), Inches(0.35))
    _text(box, [f"AZ-Eval   ·   CMAI 2026   ·   {number}"], 10, MUTED)


def _title_slide(prs) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint(slide, prs)

    _rule(slide, _margin(), Inches(1.55), Inches(2.4), height_pt=6)

    box = slide.shapes.add_textbox(_margin(), Inches(1.95), Inches(11.6), Inches(2.4))
    _text(box, TITLE.split("\n"), 34, INK, bold=True, space_after=4)

    box = slide.shapes.add_textbox(_margin(), Inches(4.45), Inches(11.6), Inches(1.0))
    _text(box, [AUTHOR, AFFILIATION], 20, INK_SOFT, space_after=2)

    box = slide.shapes.add_textbox(_margin(), Inches(6.5), Inches(11.6), Inches(0.5))
    _text(box, [VENUE], 14, MUTED)


def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9, konfrans proyektorlarının standartı
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    for number, item in enumerate(MAIN + BACKUP, start=2):
        # Tamamilə boş düzüm. Yer tutuculu düzümlərdə şrift və mövqe şablondan
        # gəlir və başqa kompüterdə açılanda dəyişə bilər; hər şeyi özümüz
        # qurmaq nəticəni proqnozlaşdırılan edir.
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _paint(slide, prs)

        box = slide.shapes.add_textbox(_margin(), Inches(0.5), Inches(11.9), Inches(0.95))
        _text(box, [item.title], 30, INK, bold=True)
        _rule(slide, _margin(), Inches(1.5), Inches(1.1))
        _footer(slide, number)

        top = Inches(1.95)

        if item.image:
            slide.shapes.add_picture(
                str(HERE / item.image), Inches(1.1), Inches(1.95), width=Inches(11.1)
            )
            if item.punch:
                box = slide.shapes.add_textbox(
                    _margin(), Inches(6.4), Inches(11.9), Inches(0.5)
                )
                _text(box, [item.punch], 16, INK_SOFT)
            _add_notes(slide, item.notes)
            continue

        # Yalnız vurğudan ibarət slayd: rəqəm və ya sual böyük yazılır və
        # slaydın mərkəzində durur. 12 dəqiqəlik çıxışda belə slayd nəfəs
        # alma nöqtəsidir.
        if item.punch and not item.bullets and not item.table:
            box = slide.shapes.add_textbox(_margin(), Inches(2.9), Inches(11.9), Inches(1.6))
            _text(box, [item.punch], 44, ACCENT, bold=True)
            _add_notes(slide, item.notes)
            continue

        if item.punch:
            # Mavi YALNIZ vurğu zolağı və tək qəhrəman cümlə üçündür. Hər
            # vurğunu mavi etsək, rəng vurğulama qabiliyyətini itirir.
            box = slide.shapes.add_textbox(_margin(), top, Inches(11.9), Inches(0.7))
            _text(box, [item.punch], 22, INK, bold=True)
            top = Inches(2.75)

        if item.table:
            headers, rows = item.table
            shape = slide.shapes.add_table(
                len(rows) + 1, len(headers), _margin(), top,
                Inches(11.9), Inches(0.42 * (len(rows) + 1)),
            )
            table = shape.table
            # Susma cədvəl üslubu mavi zolaqlıdır və slaydın sakit tonunu pozur.
            table.first_row = False
            table.horz_banding = False
            for column, text in enumerate(headers):
                _set_cell(table.cell(0, column), text, bold=True, size=17)
            for row_index, row in enumerate(rows, start=1):
                for column, text in enumerate(row):
                    _set_cell(table.cell(row_index, column), text, size=17)
            top = Inches(top.inches + 0.5 * (len(rows) + 1) + 0.35)

        if item.bullets:
            box = slide.shapes.add_textbox(
                _margin(), top, Inches(11.9), Inches(6.8 - top.inches)
            )
            _text(box, item.bullets, 19, INK_SOFT)

        _add_notes(slide, item.notes)

    prs.save(path)


def _set_cell(cell, text: str, bold: bool = False, size: int = 17) -> None:
    """Xanaya mətn qoyur və şrifti təyin edir.

    BOŞ MƏTN TƏLƏSİ: `cell.text = ""` heç bir `run` yaratmır, ona görə
    `runs[0]` IndexError verir. Bəzi cədvəllərdə başlıq sətri qəsdən boşdur
    (AZ-Eval slaydı iki sütunlu sadə siyahıdır), ona görə yoxlama lazımdır.
    """
    from pptx.util import Pt

    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(SURFACE)
    cell.text = text
    paragraph = cell.text_frame.paragraphs[0]
    if not paragraph.runs:
        return
    run = paragraph.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = _rgb(INK if bold else INK_SOFT)


def _add_notes(slide, notes: str) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# --------------------------------------------------------------------------
# Beamer
# --------------------------------------------------------------------------


def _tex_escape(text: str) -> str:
    """LaTeX-də xüsusi mənası olan simvolları qoruyur.

    Slayd mətnində `%` (faiz!), `_` (model adlarında), `&` və `#` var - qorunmasa
    kompilyasiya səssizcə pozulur və ya sətir itir.
    """
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    return text


def _beamer_frame(item: Slide) -> str:
    lines = [f"\\begin{{frame}}{{{_tex_escape(item.title)}}}"]

    if item.image:
        # 0.92 eni vurğu sətrini kadrdan çıxarırdı. Qrafiklərin özündə onsuz da
        # kənar boşluq var, ona görə kiçiltmək heç nə itirmir.
        lines.append("\\begin{center}")
        lines.append(
            f"\\includegraphics[width=0.80\\textwidth]{{{Path(item.image).stem}}}"
        )
        lines.append("\\end{center}")
        if item.punch:
            lines.append("\\vspace{-0.15cm}")
            lines.append(
                f"\\begin{{center}}{{\\small\\color{{inksoft}} {_tex_escape(item.punch)}}}\\end{{center}}"
            )
        if item.notes:
            lines.append(f"\\note{{{_tex_escape(item.notes)}}}")
        lines.append("\\end{frame}")
        return "\n".join(lines)

    if item.punch and not item.bullets and not item.table:
        lines.append("\\vspace{1.2cm}")
        lines.append("\\begin{center}")
        lines.append(
            f"{{\\Huge\\bfseries\\color{{accent}} {_tex_escape(item.punch)}}}"
        )
        lines.append("\\end{center}")
    else:
        if item.punch:
            lines.append(
                f"\\begin{{center}}{{\\large\\bfseries {_tex_escape(item.punch)}}}\\end{{center}}"
            )
            lines.append("\\vspace{0.3cm}")

        if item.table:
            headers, rows = item.table
            spec = "l" + "r" * (len(headers) - 1)
            lines.append("\\begin{center}")
            lines.append(f"\\begin{{tabular}}{{{spec}}}")
            lines.append("\\toprule")
            if any(headers):
                lines.append(
                    " & ".join(f"\\textbf{{{_tex_escape(h)}}}" for h in headers) + " \\\\"
                )
                lines.append("\\midrule")
            for row in rows:
                lines.append(" & ".join(_tex_escape(c) for c in row) + " \\\\")
            lines.append("\\bottomrule")
            lines.append("\\end{tabular}")
            lines.append("\\end{center}")
            lines.append("\\vspace{0.3cm}")

        if item.bullets:
            lines.append("\\begin{itemize}")
            for line in item.bullets:
                if not line.strip():
                    lines.append("\\vspace{0.25cm}")
                    continue
                lines.append(f"  \\item[] {_tex_escape(line)}")
            lines.append("\\end{itemize}")

    if item.notes:
        lines.append(f"\\note{{{_tex_escape(item.notes)}}}")
    lines.append("\\end{frame}")
    return "\n".join(lines)


def build_beamer(path: Path) -> None:
    frames = "\n\n".join(_beamer_frame(item) for item in MAIN)
    backup = "\n\n".join(_beamer_frame(item) for item in BACKUP)
    # Əl ilə sətir kəsimi (`\\`) QOYULMUR. Beamer-də başlıq onsuz da sətirə
    # sığmayıb keçir; hər ikisi olanda aralıqlar qeyri-bərabər çıxır. Təbii
    # keçidə buraxmaq bərabər sətir aralığı verir.
    title_tex = _tex_escape(TITLE).replace("\n", " ")

    document = f"""% CMAI 2026 - `python talk/slides.py --beamer` ilə qurulub. ƏL İLƏ REDAKTƏ ETMƏ:
% məzmun `talk/slides.py` faylındadır, bu fayl hər qurulmada üstündən yazılır.
%
% XELATEX İLƏ KOMPİLYASİYA ET, pdflatex ilə YOX:
%     xelatex slides.tex
% Səbəb: nümunələrdə `ə` (U+0259) var və standart pdflatex şriftlərində yoxdur.
\\documentclass[aspectratio=169,12pt]{{beamer}}
\\usetheme{{default}}
\\usepackage{{fontspec}}
\\usepackage{{booktabs}}
\\usepackage{{graphicx}}
\\usepackage{{xcolor}}
\\setmainfont{{{BEAMER_FONT}}}
\\setsansfont{{{BEAMER_FONT}}}

% Rənglər `talk/figures.py`-dakı təsdiqlənmiş palitra ilə eynidir, yəni
% slaydla qrafik proyektorda eyni mavini göstərir.
\\definecolor{{accent}}{{HTML}}{{{ACCENT}}}
\\definecolor{{ink}}{{HTML}}{{{INK}}}
\\definecolor{{inksoft}}{{HTML}}{{{INK_SOFT}}}
\\definecolor{{muted}}{{HTML}}{{{MUTED}}}
\\definecolor{{surface}}{{HTML}}{{{SURFACE}}}

% Fon PPTX versiyası ilə eyni isti ağdır, xalis ağ deyil.
\\setbeamercolor{{background canvas}}{{bg=surface}}

\\setbeamercolor{{frametitle}}{{fg=ink,bg=}}
\\setbeamercolor{{title}}{{fg=ink}}
\\setbeamercolor{{author}}{{fg=inksoft}}
\\setbeamercolor{{institute}}{{fg=inksoft}}
\\setbeamercolor{{date}}{{fg=inksoft}}
\\setbeamercolor{{normal text}}{{fg=ink}}
\\setbeamerfont{{frametitle}}{{size=\\large,series=\\bfseries}}

% Başlığın altında vurğu zolağı, PPTX versiyası ilə eyni.
\\setbeamertemplate{{frametitle}}{{%
  \\vspace{{0.45cm}}%
  \\usebeamerfont{{frametitle}}\\usebeamercolor[fg]{{frametitle}}\\insertframetitle\\par
  \\vspace{{0.12cm}}%
  \\hspace*{{0.02\\textwidth}}\\textcolor{{accent}}{{\\rule{{0.10\\textwidth}}{{2.2pt}}}}%
}}

\\setbeamertemplate{{navigation symbols}}{{}}
\\setbeamertemplate{{footline}}[frame number]
\\setbeamercolor{{footline}}{{fg=inksoft}}

\\title{{{title_tex}}}
\\author{{{_tex_escape(AUTHOR)}}}
\\institute{{{_tex_escape(AFFILIATION)}}}
\\date{{{_tex_escape(VENUE)}}}

\\begin{{document}}

% Başlıq səhifəsi əl ilə qurulur. Beamer-in `\\titlepage` komandası mətni
% mərkəzə yığır və PPTX versiyası ilə uyuşmur; sola düzləndirmə və vurğu
% zolağı iki formatı eyni görkəmə gətirir.
\\begin{{frame}}[plain]
  \\vspace{{1.4cm}}
  \\hspace*{{0.02\\textwidth}}\\textcolor{{accent}}{{\\rule{{0.22\\textwidth}}{{4pt}}}}\\par
  \\vspace{{0.55cm}}
  \\hspace*{{0.02\\textwidth}}\\parbox{{0.95\\textwidth}}{{\\raggedright
    {{\\LARGE\\bfseries\\color{{ink}} {title_tex}}}\\par
    \\vspace{{0.85cm}}
    {{\\large\\color{{inksoft}} {_tex_escape(AUTHOR)}}}\\par
    \\vspace{{0.12cm}}
    {{\\color{{inksoft}} {_tex_escape(AFFILIATION)}}}\\par
    \\vspace{{1.1cm}}
    {{\\small\\color{{muted}} {_tex_escape(VENUE)}}}
  }}
\\end{{frame}}

{frames}

\\appendix
\\begin{{frame}}
  \\begin{{center}}{{\\Large Backup slides}}\\end{{center}}
\\end{{frame}}

{backup}

\\end{{document}}
"""
    path.write_text(document, encoding="utf-8")


# --------------------------------------------------------------------------
# CLI
# --------------------------------------------------------------------------


def main(argv: list[str] | None = None) -> int:
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8", errors="replace")

    parser = argparse.ArgumentParser(description="CMAI 2026 slaydlarını qurur")
    parser.add_argument("--pptx", action="store_true")
    parser.add_argument("--beamer", action="store_true")
    parser.add_argument(
        "--compile",
        action="store_true",
        help="beamer .tex faylını xelatex ilə PDF-ə çevir",
    )
    args = parser.parse_args(argv)

    both = not (args.pptx or args.beamer)

    if args.pptx or both:
        target = HERE / "CMAI_2026_Garibli.pptx"
        build_pptx(target)
        print(f"PPTX  -> {target}")

    if args.beamer or both:
        target = HERE / "slides.tex"
        build_beamer(target)
        print(f"TeX   -> {target}")
        if args.compile:
            result = subprocess.run(
                ["xelatex", "-interaction=nonstopmode", target.name],
                cwd=HERE,
                capture_output=True,
                text=True,
            )
            pdf = HERE / "slides.pdf"
            if pdf.exists():
                print(f"PDF   -> {pdf}")
            else:
                print("xelatex uğursuz oldu:", file=sys.stderr)
                print(result.stdout[-2000:], file=sys.stderr)
                return 1

    print(f"\n{len(MAIN)} əsas + {len(BACKUP)} ehtiyat slayd")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
